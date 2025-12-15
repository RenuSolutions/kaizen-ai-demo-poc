import io
import re
from io import BytesIO
from typing import Dict, List

import streamlit as st
from pptx import Presentation

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn

from openai import OpenAI
from openai import AuthenticationError, RateLimitError, APIConnectionError, APIStatusError


# -----------------------------
# Page config
# -----------------------------
st.set_page_config(
    page_title="Kaizen Deck – Executive Summary (Communication-Ready)",
    layout="wide",
)

st.title("Kaizen Deck – Executive Summary (Communication-Ready)")
st.caption(
    "Upload a Kaizen report-out PPTX and generate a PPC-formatted, executive-ready one-page summary "
    "using the approved Word template placeholders."
)

# -----------------------------
# Secrets / client setup
# -----------------------------
api_key = st.secrets.get("OPENAI_API_KEY", "").strip()
if not api_key:
    st.warning(
        "⚠️ Missing `OPENAI_API_KEY`.\n\n"
        "Add it in **Manage app → Settings → Secrets** as:\n\n"
        '`OPENAI_API_KEY="sk-..."`'
    )
    st.stop()

client = OpenAI(api_key=api_key)


# -----------------------------
# Helpers
# -----------------------------
def extract_slide_text(pptx_bytes: bytes) -> str:
    """Extract plain text from all slides."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    parts.append(t)
        if parts:
            out.append(f"Slide {i}:\n" + "\n".join(parts))
    return "\n\n".join(out)


def truncate_text(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n\n[TRUNCATED FOR DEMO COST CONTROL]"


def _insert_paragraph_after(paragraph, text: str = "", style: str = None):
    """
    Insert a new paragraph AFTER 'paragraph' (python-docx doesn't provide insert_paragraph_after).
    Returns the new paragraph.
    """
    new_p = OxmlElement("w:p")
    paragraph._p.addnext(new_p)
    new_para = paragraph._parent.add_paragraph()
    # Move the newly created paragraph element into the right spot
    new_para._p = new_p

    if style:
        try:
            new_para.style = style
        except Exception:
            pass

    if text:
        new_para.add_run(text)

    return new_para


def _clean_bullets(lines: List[str]) -> List[str]:
    cleaned = []
    for ln in lines:
        ln = ln.strip()
        ln = re.sub(r"^[-•\u2022]\s*", "", ln)  # strip leading bullets
        if ln:
            cleaned.append(ln)
    return cleaned


def _parse_sections(ai_text: str) -> Dict[str, str]:
    """
    Expect AI to return tagged sections:
    [OVERVIEW]...[/OVERVIEW]
    [CHALLENGES]...[/CHALLENGES]
    etc.
    """
    tags = ["OVERVIEW", "CHALLENGES", "IMPROVEMENTS", "BENEFITS", "PLAN", "SUMMARY"]
    out = {}
    for t in tags:
        m = re.search(rf"\[{t}\](.*?)\[/\s*{t}\]", ai_text, flags=re.DOTALL | re.IGNORECASE)
        out[t] = (m.group(1).strip() if m else "").strip()
    return out


def generate_exec_summary_sections(slide_text: str) -> Dict[str, str]:
    """
    Generate content for the PPC one-page template using strict tags so we can place content reliably.
    """
    prompt = f"""
You are writing a PPC Partners executive one-page Kaizen summary.
You MUST output ONLY these six tagged sections, in this exact order, with no extra text:

[OVERVIEW]
2–3 sentences. No bullets.
[/OVERVIEW]

[CHALLENGES]
3–6 bullets max. Each bullet on its own line. No numbering.
[/CHALLENGES]

[IMPROVEMENTS]
3–6 bullets max. Each bullet on its own line. No numbering.
[/IMPROVEMENTS]

[BENEFITS]
3–6 bullets max. Each bullet on its own line. No numbering.
[/BENEFITS]

[PLAN]
3 bullets max. Each bullet on its own line. Use time horizons like 0–30, 30–90, 6–12 months if possible.
[/PLAN]

[SUMMARY]
1–2 sentences. No bullets.
[/SUMMARY]

Rules:
- Use only facts supported by the slide content; if something is missing, write "TBD" briefly.
- Keep it concise and executive-ready.
- Do not add any other headings or commentary.

Kaizen slide content:
{slide_text}
""".strip()

    resp = client.responses.create(
        model="gpt-4o-mini",
        input=prompt,
    )
    sections = _parse_sections(resp.output_text)
    return sections


def fill_docx_template(template_bytes: bytes, sections: Dict[str, str]) -> bytes:
    """
    Replace placeholders in the template with content.
    Placeholders must be on their own line:
    {{OVERVIEW}}, {{CHALLENGES}}, {{IMPROVEMENTS}}, {{BENEFITS}}, {{PLAN}}, {{SUMMARY}}
    """
    doc = Document(BytesIO(template_bytes))

    placeholders = {
        "{{OVERVIEW}}": "OVERVIEW",
        "{{CHALLENGES}}": "CHALLENGES",
        "{{IMPROVEMENTS}}": "IMPROVEMENTS",
        "{{BENEFITS}}": "BENEFITS",
        "{{PLAN}}": "PLAN",
        "{{SUMMARY}}": "SUMMARY",
    }

    # Validate placeholders exist somewhere (paragraphs)
    all_text = "\n".join([p.text.strip() for p in doc.paragraphs])
    missing = [ph for ph in placeholders.keys() if ph not in all_text]
    if missing:
        raise ValueError(
            "Template must include placeholders on their own lines: "
            "{{OVERVIEW}}, {{CHALLENGES}}, {{IMPROVEMENTS}}, {{BENEFITS}}, {{PLAN}}, {{SUMMARY}}"
        )

    for p in list(doc.paragraphs):
        key = p.text.strip()
        if key in placeholders:
            section_name = placeholders[key]
            content = sections.get(section_name, "").strip()

            # Clear the placeholder paragraph text
            p.text = ""

            if section_name in ("OVERVIEW", "SUMMARY"):
                # Single paragraph (no bullets)
                if not content:
                    content = "TBD"
                p.add_run(content)

            else:
                # Bullet list section
                lines = _clean_bullets(content.splitlines())
                if not lines:
                    lines = ["TBD"]

                # Put first bullet in the same paragraph (so we don’t leave blank space)
                try:
                    p.style = "List Bullet"
                except Exception:
                    pass
                p.add_run(lines[0])

                # Remaining bullets inserted after
                prev = p
                for ln in lines[1:]:
                    newp = _insert_paragraph_after(prev, "", style="List Bullet")
                    newp.add_run(ln)
                    prev = newp

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# -----------------------------
# UI
# -----------------------------
col1, col2 = st.columns(2)

with col1:
    pptx_file = st.file_uploader("Upload Kaizen Report-Out Deck (PPTX)", type=["pptx"])

with col2:
    template_file = st.file_uploader("Upload PPC Executive Summary Word Template (.docx)", type=["docx"])

st.info("Upload both the PPTX and the approved Word template to continue.")

if not pptx_file or not template_file:
    st.stop()

pptx_bytes = pptx_file.read()
template_bytes = template_file.read()

st.success("Files uploaded successfully.")

# Cost control
st.subheader("Cost Controls")
max_chars = st.slider(
    "Max characters of slide text sent to the model",
    min_value=5_000,
    max_value=60_000,
    value=20_000,
    step=5_000,
)

with st.expander("Preview extracted slide text (truncated)"):
    slide_text_full = extract_slide_text(pptx_bytes)
    slide_text = truncate_text(slide_text_full, max_chars=max_chars)
    st.text(slide_text[:12000])

st.divider()

if st.button("Generate Executive Summary"):
    try:
        with st.spinner("Generating summary content..."):
            slide_text_full = extract_slide_text(pptx_bytes)
            slide_text = truncate_text(slide_text_full, max_chars=max_chars)
            sections = generate_exec_summary_sections(slide_text)

        with st.spinner("Filling Word template..."):
            output_docx = fill_docx_template(template_bytes, sections)

        st.success("Executive Summary generated.")

        st.download_button(
            label="Download Executive Summary (.docx)",
            data=output_docx,
            file_name="Executive_Summary_One_Page.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

        st.subheader("Preview (what was inserted)")
        st.write({"OVERVIEW": sections["OVERVIEW"]})
        st.write({"CHALLENGES": sections["CHALLENGES"]})
        st.write({"IMPROVEMENTS": sections["IMPROVEMENTS"]})
        st.write({"BENEFITS": sections["BENEFITS"]})
        st.write({"PLAN": sections["PLAN"]})
        st.write({"SUMMARY": sections["SUMMARY"]})

    except ValueError as ve:
        st.error(str(ve))
    except AuthenticationError:
        st.error("OpenAI authentication failed. Check your OPENAI_API_KEY in Streamlit Secrets.")
    except RateLimitError:
        st.error("OpenAI quota/rate limit hit. Check billing/usage limits, then try again.")
    except APIConnectionError:
        st.error("Network/API connection error. Try again.")
    except APIStatusError as e:
        st.error(f"OpenAI API returned an error: {e}")
    except Exception as e:
        st.error(f"Unexpected error: {e}")


